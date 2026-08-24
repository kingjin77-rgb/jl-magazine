#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""편집 가능한 PPTX 내보내기 — 배경은 이미지, 글자는 텍스트 상자.

같은 슬라이드를 두 번 렌더한다.

  1) 글자를 투명하게 만든 **배경 PNG** — 도해·표 괘선·밴드·사진만 남는다
  2) 같은 페이지에서 **글자의 위치·크기·색·굵기를 JSON으로 덤프**

PPTX는 배경 PNG 위에 덤프대로 텍스트 상자를 얹는다. PPT에서 문구를 고칠 수 있다.

  python3 tools/export_pptx.py
  python3 tools/export_pptx.py --check              # 검증 오버레이까지
  python3 tools/export_pptx.py --only S32_Closing   # 한 장만
"""
import io, json, os, sys
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BUILD = os.path.join(ROOT, 'template', 'design', 'build')
OUT_DIR = os.path.join(ROOT, 'template', 'design', 'preview')
BG_DIR = os.path.join(OUT_DIR, '_bg')
OUT = os.path.join(OUT_DIR, 'JL_하자소송_기본제안서.pptx')
CHROME = '/opt/pw-browsers/chromium-1194/chrome-linux/chrome'

W_PX, H_PX = 1332, 750
EMU_PER_PX = 12192000 / W_PX
PT_PER_PX = 0.72

HIDE_TEXT = """
  *{color:transparent!important;-webkit-text-fill-color:transparent!important;}
  svg text,svg tspan{fill:transparent!important;}
"""

# 브라우저 안에서 도는 수집기. 텍스트를 가진 블록마다 위치·스타일·런을 뽑는다.
COLLECT = r"""
() => {
  const INLINE = new Set(['B','STRONG','EM','I','SPAN','SMALL','U','A','SUB','SUP','TSPAN','BR']);
  // 태그가 인라인이어도 display를 블록으로 바꿔 둔 경우가 있다. 계산된 값으로 판단한다.
  const isInline = k => k.tagName === 'BR' ||
    (INLINE.has(k.tagName) && getComputedStyle(k).display.startsWith('inline'));
  const out = [];

  const rgb = c => {
    const m = (c || '').match(/rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)/);
    if (!m) return null;
    if (m[4] !== undefined && parseFloat(m[4]) < 0.06) return null;
    return [ +m[1], +m[2], +m[3] ];
  };
  const fam = cs => (cs.fontFamily.split(',')[0] || '').replace(/["']/g, '');

  const runsOf = (el, base) => {
    const runs = [];
    for (const n of el.childNodes) {
      if (n.nodeType === 3) {
        const t = n.textContent.replace(/\s+/g, ' ');
        if (t.trim()) runs.push({ text: t, bold: base.bold, color: base.color,
                                  size: base.size, font: base.font });
      } else if (n.nodeType === 1) {
        if (n.tagName === 'BR') { runs.push({ br: true }); continue; }
        const cs = getComputedStyle(n);
        const t = n.textContent.replace(/\s+/g, ' ');
        if (!t.trim()) continue;
        runs.push({ text: t,
                    bold: parseInt(cs.fontWeight, 10) >= 600,
                    color: rgb(cs.color) || base.color,
                    size: parseFloat(cs.fontSize) || base.size,
                    font: fam(cs) });
      }
    }
    return runs;
  };

  const walk = el => {
    const cs = getComputedStyle(el);
    if (cs.visibility === 'hidden' || cs.display === 'none' || parseFloat(cs.opacity) === 0) return;

    const kids = [...el.children];
    const hasText = [...el.childNodes].some(n => n.nodeType === 3 && n.textContent.trim());
    const isSvgText = el.tagName === 'text';
    const blockish = isSvgText || (hasText && kids.every(isInline));

    const base = { bold: parseInt(cs.fontWeight, 10) >= 600,
                   color: rgb(isSvgText ? cs.fill : cs.color) || [17, 20, 24],
                   size: parseFloat(cs.fontSize) || 14,
                   font: fam(cs) };
    const pad = k => parseFloat(cs['padding' + k]) || 0;
    const lh = parseFloat(cs.lineHeight) || (base.size * 1.4);
    const align = isSvgText ? ({ middle: 'center', end: 'right' }[cs.textAnchor] || 'left')
                            : cs.textAlign;

    // 텍스트와 블록 자식이 섞인 경우: 자기 텍스트만 먼저 뽑고 자식으로 내려간다
    if (hasText && !blockish) {
      const r = el.getBoundingClientRect();
      const own = [];
      for (const n of el.childNodes) {
        if (n.nodeType !== 3) continue;
        const t = n.textContent.replace(/\s+/g, ' ');
        if (t.trim()) own.push({ text: t, bold: base.bold, color: base.color,
                                 size: base.size, font: base.font });
      }
      if (own.length && r.width > 1) {
        out.push({ x: r.left, y: r.top, w: r.width, h: Math.min(r.height, lh * 2),
                   pl: pad('Left'), pr: pad('Right'), pt: pad('Top'), pb: 0,
                   align: align, lh: lh, size: base.size, font: base.font, runs: own });
      }
      for (const k of kids) walk(k);
      return;
    }

    if (blockish) {
      const r = el.getBoundingClientRect();
      if (r.width > 1 && r.height > 1 && r.top < 760 && r.left < 1340) {
        const runs = isSvgText
          ? [{ text: el.textContent.replace(/\s+/g, ' '), ...base }]
          : runsOf(el, base);
        if (runs.length) {
          out.push({ x: r.left, y: r.top, w: r.width, h: r.height,
                     pl: pad('Left'), pr: pad('Right'), pt: pad('Top'), pb: pad('Bottom'),
                     align: align, lh: lh, size: base.size, font: base.font, runs: runs });
        }
      }
      return;                       // 블록 안으로 더 들어가지 않는다
    }
    for (const k of kids) walk(k);
    if (el.tagName === 'svg') for (const t of el.querySelectorAll('text')) walk(t);
  };

  walk(document.body);
  return out;
}
"""


def artboards():
    m = json.load(io.open(os.path.join(BUILD, 'canvas.json'), encoding='utf-8'))
    return [os.path.join(BUILD, a['file']) for a in m['artboards']
            if os.path.exists(os.path.join(BUILD, a['file']))]


def harvest(only=None):
    """배경 PNG와 텍스트 덤프를 한 번의 브라우저 세션에서 만든다."""
    os.makedirs(BG_DIR, exist_ok=True)
    files = artboards()
    if only:
        files = [f for f in files if only in os.path.basename(f)]
    result = []
    with sync_playwright() as pw:
        br = pw.chromium.launch(executable_path=CHROME, args=['--no-sandbox'])
        pg = br.new_page(viewport={'width': W_PX, 'height': H_PX})
        for i, f in enumerate(files):
            name = os.path.basename(f).replace('.dc.html', '')
            pg.goto('file://' + f)
            pg.wait_for_timeout(320)
            texts = pg.evaluate(COLLECT)
            pg.add_style_tag(content=HIDE_TEXT)          # 글자만 지운 배경
            pg.wait_for_timeout(90)
            bg = os.path.join(BG_DIR, '%02d_%s.png' % (i + 1, name))
            pg.screenshot(path=bg)
            result.append((name, bg, texts))
            print('  %02d %-28s 텍스트 %d블록' % (i + 1, name, len(texts)))
        br.close()
    return result


ALIGN = {'left': PP_ALIGN.LEFT, 'start': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
         'right': PP_ALIGN.RIGHT, 'end': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY}
FONT = {'JL디스플레이': '210 옴니고딕 050', 'JL본문': 'Pretendard'}


def geom(t):
    """PPT 텍스트 상자의 위치·크기. 폰트 메트릭 차이로 줄이 깨지지 않게 여유를 준다.
    우측 정렬 블록은 오른쪽 끝이 기준이므로 왼쪽으로 넓힌다."""
    one = t['h'] <= t['lh'] * 1.6
    w = max(t['w'], 8) * (1.28 if one else 1.05) + 14
    x = t['x'] - (w - t['w']) if t['align'] in ('right', 'end') else t['x']
    h = max(t['h'], 8) * (1.5 if one else 1.08)
    return max(x, 0), t['y'], w, h


def build(items):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for name, bg, texts in items:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_picture(bg, 0, 0, width=Emu(12192000), height=Emu(6858000))
        for t in texts:
            x, y, w, h = geom(t)
            box = s.shapes.add_textbox(Emu(int(x * EMU_PER_PX)), Emu(int(y * EMU_PER_PX)),
                                       Emu(int(w * EMU_PER_PX)), Emu(int(h * EMU_PER_PX)))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = Emu(int(t['pl'] * EMU_PER_PX))
            tf.margin_right = Emu(int(t['pr'] * EMU_PER_PX))
            tf.margin_top = Emu(int(t['pt'] * EMU_PER_PX))
            tf.margin_bottom = Emu(int(t['pb'] * EMU_PER_PX))
            al = ALIGN.get(t['align'], PP_ALIGN.LEFT)
            para = tf.paragraphs[0]
            para.alignment = al
            para.line_spacing = Pt(t['lh'] * PT_PER_PX)
            for r in t['runs']:
                if r.get('br'):
                    para = tf.add_paragraph()
                    para.alignment = al
                    para.line_spacing = Pt(t['lh'] * PT_PER_PX)
                    continue
                run = para.add_run()
                run.text = r['text']
                f = run.font
                f.size = Pt(round(r.get('size', t['size']) * PT_PER_PX, 1))
                f.bold = bool(r.get('bold'))
                f.color.rgb = RGBColor(*(r.get('color') or [17, 20, 24]))
                fam = r.get('font') or t.get('font') or 'Pretendard'
                f.name = FONT.get(fam, fam)
        s.notes_slide.notes_text_frame.text = name
    prs.save(OUT)
    print('PPTX %s · %.1f MB · %d장' % (OUT, os.path.getsize(OUT) / 1024 / 1024, len(items)))


def overlay(items):
    """검증용 — 덤프한 글자를 배경 PNG 위에 PPT와 같은 좌표·폭으로 다시 그린다.

    이 컨테이너의 리브레오피스가 pptx를 열지 못해(python-pptx 최소 파일도 실패)
    PPTX를 직접 렌더할 수 없다. 대신 PPT가 받게 될 것과 같은 조건으로 재현해 확인한다."""
    from PIL import Image, ImageDraw, ImageFont
    FDIR = os.path.join(ROOT, 'assets', 'fonts')
    reg, bold = os.path.join(FDIR, 'Pretendard-Medium.ttf'), os.path.join(FDIR, 'Pretendard-SemiBold.ttf')
    cache = {}

    def font(px, b):
        k = (int(px), bool(b))
        if k not in cache:
            cache[k] = ImageFont.truetype(bold if b else reg, max(7, int(px)))
        return cache[k]

    for name, bg, texts in items:
        im = Image.open(bg).convert('RGB')
        d = ImageDraw.Draw(im)
        for t in texts:
            x0, y0, w, _ = geom(t)
            x = x0 + t['pl']; y = y0 + t['pt']; maxw = w - t['pl'] - t['pr']
            cx = x
            for r in t['runs']:
                if r.get('br'):
                    cx = x; y += t['lh']; continue
                sz = r.get('size', t['size'])
                f = font(sz, r.get('bold'))
                for word in r['text'].split(' '):
                    ww = d.textlength(word + ' ', font=f)
                    if cx + ww > x + maxw + 1 and cx > x:
                        cx = x; y += t['lh']
                    d.text((cx, y + (t['lh'] - sz) * .45), word + ' ', font=f,
                           fill=tuple(r.get('color') or (17, 20, 24)))
                    cx += ww
        im.save(os.path.join(BG_DIR, '_check_%s.png' % name))
    print('검증 오버레이 %d장 → %s' % (len(items), BG_DIR))


if __name__ == '__main__':
    only = sys.argv[sys.argv.index('--only') + 1] if '--only' in sys.argv else None
    items = harvest(only)
    build(items)
    if '--check' in sys.argv:
        overlay(items)
